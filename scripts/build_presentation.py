# -*- coding: utf-8 -*-
"""2차 발표 슬라이드(.pptx) 생성 — 리뉴얼 디자인 (Pretendard · 타이포 조판 · 헤어라인).

근거: docs/발표_PPT_구성안.md (3절 13장 스펙), docs/발표_대사.md
디자인 원칙: 카드 최소화, 비대칭 그리드, 강조 슬라이드당 1회, 표는 헤어라인만, 이모지 없음.
전제: 시스템에 Pretendard 설치 필요. 발표는 본인 PC 또는 폰트 임베드 PDF로.

사용법: python scripts/build_presentation.py <assets_dir>
  assets_dir 에 slide6@2x.png, arch@2x.png, slide11@2x.png (docs/*.svg → Chrome headless 변환)
출력: docs/발표_슬라이드.pptx
"""
import sys
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 색 (브랜드 유지) ----
CREAM = RGBColor(0xF6, 0xF5, 0xF0)
PAPER = RGBColor(0xFB, 0xFB, 0xF8)
FOREST = RGBColor(0x13, 0x30, 0x1C)
FOREST2 = RGBColor(0x1C, 0x3E, 0x28)
GREEN = RGBColor(0x2E, 0x9E, 0x57)
GREEN_DK = RGBColor(0x1F, 0x80, 0x47)
SAGE = RGBColor(0xA8, 0xC6, 0x6C)
SAGE_DK = RGBColor(0x6F, 0x8F, 0x3A)
INK = RGBColor(0x1B, 0x24, 0x1D)
MUTED = RGBColor(0x66, 0x70, 0x6A)
FAINT = RGBColor(0x9A, 0xA4, 0x9B)
HAIR = RGBColor(0xD8, 0xDC, 0xD2)
HAIR_DK = RGBColor(0x2C, 0x47, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xC9, 0xD6, 0xCC)
AMBER = RGBColor(0xC2, 0x6B, 0x2E)
TINT = RGBColor(0xEC, 0xF5, 0xEE)

# ---- Pretendard weights ----
BLACK = "Pretendard Black"
XBOLD = "Pretendard ExtraBold"
SEMI = "Pretendard SemiBold"
MED = "Pretendard Medium"
REG = "Pretendard"
LIGHT = "Pretendard Light"

SW, SH = Inches(13.333), Inches(7.5)
LM = Inches(0.85)          # 좌 여백
RM = Inches(12.48)         # 우 경계
CW = Inches(11.63)         # 콘텐츠 폭


def _run(p, t, size, font, color, spc=None):
    r = p.add_run()
    r.text = t
    r.font.size = Pt(size)
    r.font.name = font
    r.font.bold = False
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {}); rPr.append(ea); ea.set("typeface", font)
    if spc is not None:
        rPr.set("spc", str(int(spc * 100)))
    return r


def text(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.08, sa=0):
    b = s.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls; p.space_after = Pt(sa)
        for spec in (line if isinstance(line, list) else [line]):
            _run(p, *spec)
    return b


def hline(s, x, y, w, color, weight=1.0, dash=None):
    ln = s.shapes.add_connector(2, x, y, Emu(int(x) + int(w)), y)
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    if dash:
        d = ln.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": dash})
        ln.line._get_or_add_ln().append(d)
    return ln


def vline(s, x, y, h, color, weight=1.0):
    ln = s.shapes.add_connector(2, x, y, x, Emu(int(y) + int(h)))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    return ln


def block(s, x, y, w, h, fill, radius=0.06):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                            x, y, w, h)
    if radius:
        try: sp.adjustments[0] = radius
        except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    sp.shadow.inherit = False
    sp.text_frame.paragraphs[0].text = ""
    return sp


def bg(s, color):
    return block(s, 0, 0, SW, SH, color, radius=0)


def kicker_head(s, kicker, head, accent=None):
    """좌상단 킥커 + 헤드라인 (카드 없음)."""
    text(s, LM, Inches(0.72), Inches(10), Inches(0.3), [[(kicker, 12, SEMI, GREEN, 2.5)]])
    runs = head if isinstance(head, list) else [(head, 33, XBOLD, INK, -1.0)]
    text(s, Emu(int(LM) - Inches(0.02)), Inches(1.06), Inches(11.8), Inches(0.9), [runs])


MARK_SIZE = Inches(0.22)


def logo_mark(s, x, y, size=MARK_SIZE, on_dark=False):
    """타이포그래피 모노그램 — 타이틀 슬라이드의 FOREST 배경 위 흰 '팜바톤'과 같은 배색을
    작은 사각 마크로 반복. 래스터 로고 대신 벡터 도형만 사용해 톤을 유지."""
    fill = WHITE if on_dark else FOREST
    fg = FOREST if on_dark else WHITE
    block(s, x, y, size, size, fill, radius=0.32)
    text(s, x, Emu(int(y) - Inches(0.008)), size, size, [[("팜", 12.5, BLACK, fg, 0)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(s, num, dark=False):
    c = FAINT if not dark else RGBColor(0x6E, 0x78, 0x6F)
    mark_y = Inches(7.03)
    logo_mark(s, LM, mark_y, on_dark=dark)
    text(s, Emu(int(LM) + int(MARK_SIZE) + Inches(0.09)), Inches(7.06), Inches(6), Inches(0.3),
         [[("팜바톤 · FarmBaton", 9.5, MED, c, 0.5)]])
    text(s, Inches(12.0), Inches(7.06), Inches(0.5), Inches(0.3),
         [[(num, 10, MED, c, 0)]], align=PP_ALIGN.RIGHT)


# ============ 표 (S3 empty / S9 filled) ============
def comparison(s, filled):
    tx = LM
    c0 = Inches(2.0)
    cw = Inches(2.55)
    fbw = Inches(2.35)
    xs = [tx, Emu(int(tx) + int(c0)), Emu(int(tx) + int(c0) + int(cw)),
          Emu(int(tx) + int(c0) + 2 * int(cw))]
    fb_x = Emu(int(tx) + int(CW) - int(fbw))
    top = Inches(2.15)
    if filled:
        block(s, Emu(int(fb_x) - Inches(0.22)), top, Emu(int(fbw) + Inches(0.22)),
              Inches(4.35), TINT, radius=0.04)
    head_y = Inches(2.32)
    heads = ["", "감정평가·중개", "부동산 플랫폼", "농지은행"]
    for i, h in enumerate(heads):
        if h:
            text(s, xs[i], head_y, cw, Inches(0.4), [[(h, 13, MED, MUTED, 0)]])
    text(s, fb_x, head_y, fbw, Inches(0.4),
         [[("팜바톤", 15, XBOLD, GREEN if filled else MUTED, 0)]])
    hline(s, tx, Inches(2.86), CW, INK, 1.5)
    rows = [
        ("평가 대상", ["토지·건물", "토지·건물(범용)", "농지만"], "농장 전체"),
        ("가치 산정", ["감정평가 · 수일", "실거래·AVM 조회", "공시지가·호가"], "결정론 산식 · 무료"),
        ("승계 매칭", ["—", "—", "농지 한정"], "6요인 양면 매칭"),
        ("당사자 연결", ["대면·전화", "매물 게시형", "매물 게시형"], "인앱 상담 → 채팅"),
    ]
    ry, rh = Inches(3.12), Inches(0.82)
    for ri, (label, vals, fb) in enumerate(rows):
        cy = Emu(int(ry) + ri * int(rh))
        text(s, xs[0], cy, c0, rh, [[(label, 14, SEMI, INK, 0)]], anchor=MSO_ANCHOR.MIDDLE)
        for ci, v in enumerate(vals):
            text(s, xs[ci + 1], cy, cw, rh,
                 [[(v, 13, REG, FAINT if v == "—" else MUTED, 0)]], anchor=MSO_ANCHOR.MIDDLE)
        fb_col = FOREST if filled else FAINT
        fb_font = SEMI if filled else REG
        text(s, fb_x, cy, fbw, rh, [[(fb if filled else "·", 13.5, fb_font, fb_col, 0)]],
             anchor=MSO_ANCHOR.MIDDLE)
        if ri < len(rows) - 1:
            hline(s, tx, Emu(int(cy) + int(rh)), CW, HAIR, 0.75)


def build(assets: Path, out: Path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]
    new = lambda: prs.slides.add_slide(blank)

    qr = assets / "qr.png"
    qrcode.make("https://farmbaton.vercel.app", box_size=10, border=0).save(qr)

    # ===== S1 타이틀 =====
    s = new(); bg(s, FOREST)
    text(s, LM, Inches(0.72), Inches(11.6), Inches(0.3),
         [[("제11회 농업·농촌 공공데이터 + AI 활용 창업경진대회", 11.5, MED, SAGE, 2.2)]])
    hline(s, Emu(int(LM) + Inches(0.02)), Inches(1.16), Inches(3.2), HAIR_DK, 1.0)
    text(s, Inches(0.8), Inches(2.3), Inches(11), Inches(2.0), [[("팜바톤", 96, BLACK, WHITE, -2.0)]])
    text(s, Inches(0.9), Inches(4.12), Inches(11), Inches(0.5), [[("FarmBaton", 22, MED, PALE, 4.0)]])
    text(s, Emu(int(LM) + Inches(0.02)), Inches(5.12), Inches(11.5), Inches(0.7),
         [[("떠나는 농장과 시작하는 청년을 ", 27, XBOLD, WHITE, -0.5), ("잇다", 27, BLACK, SAGE, -0.5)]])
    hline(s, Emu(int(LM) + Inches(0.02)), Inches(6.5), Inches(11.6), HAIR_DK, 1.0)
    logo_mark(s, LM, Inches(6.69), on_dark=True)
    text(s, Emu(int(LM) + int(MARK_SIZE) + Inches(0.1)), Inches(6.72), Inches(8), Inches(0.4),
         [[("farmbaton.vercel.app", 15, SEMI, WHITE, 0), ("   지금 접속하면 실제로 동작합니다", 12.5, REG, RGBColor(0x8F, 0xA5, 0x95), 0)]])
    text(s, Inches(9.5), Inches(6.72), Inches(2.95), Inches(0.4),
         [[("김정훈 · 윤채원 · 유수민", 11.5, MED, RGBColor(0x8F, 0xA5, 0x95), 0.5)]], align=PP_ALIGN.RIGHT)

    # ===== S2 빅넘버 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "주제 시의성", "농장은 멈추는데, 이어받을 통로가 없다")
    cols = [("51.3", "%", "65세 이상 농가인구 비율", "128만 6천 명 · 전체 인구의 2.5배", FOREST),
            ("54.8", "만 가구", "70세 이상 경영주", "전체 경영주의 44.1%", FOREST),
            ("1.1", "%", "40세 미만 청년 경영주", "1만 4천 가구 · 100명 중 1명", AMBER)]
    x0, colw, gap, top = LM, Inches(3.85), Inches(0.15), Inches(2.7)
    for i, (big, unit, lbl, sub, col) in enumerate(cols):
        x = Emu(int(x0) + i * (int(colw) + int(gap)))
        if i > 0:
            vline(s, Emu(int(x) - int(gap) // 2), Inches(2.82), Inches(2.35), HAIR, 1.0)
        text(s, x, top, colw, Inches(1.3), [[(big, 74, BLACK, col, -2.0), (unit, 22, MED, MUTED, 0)]])
        hline(s, Emu(int(x) + Inches(0.05)), Inches(4.28), Inches(2.4), col, 2.5)
        text(s, x, Inches(4.45), colw, Inches(0.9),
             [[(lbl, 15, SEMI, INK, 0)], [(sub, 11.5, REG, MUTED, 0)]], ls=1.25, sa=2)
    hline(s, LM, Inches(6.12), CW, HAIR, 1.0)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.36), Inches(11.6), Inches(0.6),
         [[("수십만 농가가 은퇴를 앞두고 있지만 — ", 19, MED, INK, -0.3),
           ("청년농과 만날 통로가 없습니다", 19, XBOLD, GREEN, -0.3)]])
    text(s, LM, Inches(6.78), Inches(8), Inches(0.3), [[("통계청 2025 농림어업총조사(잠정) · 2026.4 발표", 10, REG, FAINT, 0.5)]])
    footer(s, "02")

    # ===== S3 공백 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "기존 대안의 공백", "기존 수단은 ‘농장’을 통째로 다루지 못한다")
    comparison(s, filled=False)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.7), Inches(11.6), Inches(0.5),
         [[("농장 = 농지 + 작목 + 시설 + 판로", 17, XBOLD, INK, -0.3),
           ("  이 전체를 평가하고 연결하는 통로가 비어 있습니다", 14, REG, MUTED, 0)]])
    footer(s, "03")

    # ===== S4 솔루션 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "솔루션", "주소 하나로 시작하는 양면 플랫폼")
    cy = Inches(2.55)
    block(s, Inches(4.92), cy, Inches(3.5), Inches(2.1), FOREST, radius=0.08)
    text(s, Inches(5.12), Inches(2.95), Inches(3.1), Inches(1.4),
         [[("팜바톤", 24, BLACK, WHITE, -0.5)], [("진단 엔진 · 매칭 엔진", 13, SEMI, SAGE, 0)],
          [("공공데이터 7종 결합", 11.5, REG, PALE, 0)]], align=PP_ALIGN.CENTER, ls=1.35, sa=4)
    text(s, Inches(0.85), Inches(3.0), Inches(3.7), Inches(1.3),
         [[("농장주", 20, XBOLD, INK, 0)], [("주소만 입력하면", 13, REG, MUTED, 0)],
          [("인수 검토 리포트", 14.5, SEMI, GREEN_DK, 0)]], align=PP_ALIGN.RIGHT, ls=1.3, sa=3)
    text(s, Inches(8.75), Inches(3.0), Inches(3.7), Inches(1.3),
         [[("청년농", 20, XBOLD, INK, 0)], [("조건만 입력하면", 13, REG, MUTED, 0)],
          [("매칭 리스트", 14.5, SEMI, GREEN_DK, 0)]], align=PP_ALIGN.LEFT, ls=1.3, sa=3)
    text(s, Inches(4.55), Inches(3.42), Inches(0.4), Inches(0.4), [[("→", 20, REG, FAINT, 0)]], align=PP_ALIGN.CENTER)
    text(s, Inches(8.35), Inches(3.42), Inches(0.4), Inches(0.4), [[("←", 20, REG, FAINT, 0)]], align=PP_ALIGN.CENTER)
    hline(s, LM, Inches(5.5), CW, HAIR, 1.0)
    text(s, LM, Inches(5.8), Inches(5.6), Inches(0.8),
         [[("①  ", 15, XBOLD, GREEN, 0), ("농가 등록 → 인수 검토 리포트", 15, SEMI, INK, 0)],
          [("주소 입력만으로 필지·가치 자동 산출", 12, REG, MUTED, 0)]], ls=1.4, sa=3)
    text(s, Inches(6.9), Inches(5.8), Inches(5.5), Inches(0.8),
         [[("②  ", 15, XBOLD, GREEN, 0), ("청년농 프로필 → 매칭 리스트", 15, SEMI, INK, 0)],
          [("6요인 점수순 추천 + 인앱 상담", 12, REG, MUTED, 0)]], ls=1.4, sa=3)
    footer(s, "04")

    # ===== S5 데모 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "라이브 데모", "설명 대신, 직접 보여드리겠습니다")
    steps = [("01", "주소 입력", "팜맵 필지 면적 자동 취득"),
             ("02", "검토 리포트", "인수 검토가 범위 + 신뢰등급"),
             ("03", "청년농 매칭", "6요인 100점 점수순 추천"),
             ("04", "상담 · 채팅", "수락 시 인앱 양방향 대화")]
    x0, colw = LM, Inches(2.85)
    ny = Inches(3.1)
    hline(s, Emu(int(LM) + Inches(0.1)), Inches(3.35), Inches(10.5), HAIR, 1.0)
    for i, (n, t, d) in enumerate(steps):
        x = Emu(int(x0) + i * int(colw))
        text(s, x, ny, Inches(2.6), Inches(0.9), [[(n, 40, BLACK, GREEN, 0)]])
        text(s, x, Inches(4.25), Inches(2.6), Inches(0.4), [[(t, 17, XBOLD, INK, -0.3)]])
        text(s, x, Inches(4.7), Inches(2.55), Inches(0.7), [[(d, 12, REG, MUTED, 0)]], ls=1.3)
    hline(s, LM, Inches(6.1), CW, HAIR, 1.0)
    text(s, LM, Inches(6.35), Inches(11), Inches(0.4),
         [[("farmbaton.vercel.app", 15, SEMI, GREEN_DK, 0),
           ("   실서비스 라이브 · 백업 영상 준비", 12.5, REG, MUTED, 0)]])
    footer(s, "05")

    # ===== S6 공공데이터 (풀블리드) =====
    s = new()
    s.shapes.add_picture(str(assets / "slide6@2x.png"), 0, 0, SW, SH)

    # ===== S7 AI =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "AI 활용", "숫자는 알고리즘, 설명은 AI")
    midx = Inches(6.66)
    vline(s, midx, Inches(2.35), Inches(2.75), HAIR, 1.0)
    text(s, LM, Inches(2.35), Inches(5.4), Inches(2.9),
         [[("결정론 엔진", 20, XBOLD, GREEN_DK, 0), ("  금액·점수", 13, MED, MUTED, 0)],
          [("", 8, REG, INK, 0)],
          [("금액 계산에 LLM을 쓰지 않습니다", 15, SEMI, INK, 0)],
          [("100% 결정론 Python 산식 · pytest 55개로 고정", 12.5, REG, MUTED, 0)],
          [("같은 입력 → 같은 결과 · 추적·검증 가능", 12.5, REG, MUTED, 0)],
          [("환각으로 인한 금액 오류 원천 차단", 12.5, REG, MUTED, 0)]], ls=1.5, sa=4)
    text(s, Inches(7.0), Inches(2.35), Inches(5.3), Inches(2.9),
         [[("생성형 AI", 20, XBOLD, SAGE_DK, 0), ("  설명 전담", 13, MED, MUTED, 0)],
          [("", 8, REG, INK, 0)],
          [("같은 수치를 관점별 언어로 번역합니다", 15, SEMI, INK, 0)],
          [("농장주: 매도 준비·시설 가치 관점", 12.5, REG, MUTED, 0)],
          [("청년농: 인수 자금·정책자금·체크리스트 관점", 12.5, REG, MUTED, 0)],
          [("농장 수 × 관점 수 — 사람이 못 쓰는 규모의 개인화", 12.5, REG, MUTED, 0)]], ls=1.5, sa=4)
    hline(s, LM, Inches(5.7), CW, HAIR, 1.0)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.05), Inches(11.8), Inches(0.8),
         [[("AI를 ", 26, XBOLD, INK, -0.5), ("어디에 안 썼는지", 26, BLACK, GREEN, -0.5),
           ("가 저희의 설계입니다", 26, XBOLD, INK, -0.5)]])
    footer(s, "07")

    # ===== S8 기술 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "기술성 · 완성도", "시제품이 아니라, 운영 중인 서비스")
    pw = Inches(6.9); ph = Emu(int(pw * 620 / 1000))
    s.shapes.add_picture(str(assets / "arch@2x.png"), LM, Inches(2.2), pw, ph)
    badges = [("pytest 55개 통과", "산식은 formula.md 예시 케이스로 TDD"),
              ("신뢰등급 A~D", "자료를 낼수록 정밀해지는 단계적 평가"),
              ("전 외부 API 정적 폴백", "데모가 죽지 않는 신뢰성 설계"),
              ("V-World 차단 → 국내 프록시", "공공데이터 국외반출 제약 실전 해결")]
    bx = Inches(8.1); bw = Inches(4.38)
    for i, (t, d) in enumerate(badges):
        y = Emu(int(Inches(2.35)) + i * int(Inches(1.02)))
        text(s, bx, y, bw, Inches(0.9),
             [[(t, 15, SEMI, INK, 0)], [(d, 11.5, REG, MUTED, 0)]], ls=1.35, sa=3)
        if i < 3:
            hline(s, bx, Emu(int(y) + int(Inches(0.82))), bw, HAIR, 0.75)
    text(s, LM, Inches(6.75), CW, Inches(0.4),
         [[("Vercel · Railway · Supabase(PostGIS) — production 운영 중 · farmbaton.vercel.app", 11.5, REG, FAINT, 0)]])
    footer(s, "08")

    # ===== S9 독창성 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "독창성", "국내에 없던 ‘승계 특화’ 가치평가")
    comparison(s, filled=True)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.7), Inches(11.7), Inches(0.5),
         [[("수령계수 × 시설 잔존가 × 영업권의 결합", 16, XBOLD, INK, -0.3),
           ("  확인된 범위에서 유사 사례를 찾지 못했습니다", 13, REG, MUTED, 0)]])
    footer(s, "09")

    # ===== S10 시장 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "발전 가능성 · 시장", "54만 가구가 기다리는 시장")
    tiers = [("54만 8천 가구", "후계 연결이 필요한 70세 이상 경영주 (44.1%)", FOREST, 44),
             ("3개 도 48개 시군", "충북·경북·충남 · 과수 3작목 (시작점)", GREEN, 32),
             ("MVP 검증", "초기 타깃 — 승계 희망 등록 농가", SAGE, 24)]
    ty = Inches(2.5)
    for i, (big, lbl, col, sz) in enumerate(tiers):
        y = Emu(int(ty) + i * int(Inches(1.25)))
        text(s, LM, y, Inches(6.4), Inches(1.1),
             [[(big, sz, BLACK, col, -1.0)], [(lbl, 12.5, REG, MUTED, 0)]], ls=1.15, sa=2)
        if i < 2:
            hline(s, LM, Emu(int(y) + int(Inches(1.08))), Inches(6.3), HAIR, 0.75)
    vline(s, Inches(7.6), Inches(2.55), Inches(3.7), HAIR, 1.0)
    text(s, Inches(8.0), Inches(2.75), Inches(4.4), Inches(0.4),
         [[("승계 검토의 진입장벽", 16, SEMI, INK, 0)]])
    text(s, Inches(8.0), Inches(3.5), Inches(4.4), Inches(0.4),
         [[("지금", 13, MED, MUTED, 0), ("   감정평가 수십만 원 · 수일", 14, REG, MUTED, 0)]])
    text(s, Inches(8.0), Inches(4.2), Inches(4.4), Inches(1.0),
         [[("팜바톤", 13, MED, GREEN_DK, 0)], [("무료 · 수십 초", 40, BLACK, GREEN, -1.0)]], ls=1.15, sa=4)
    text(s, Inches(8.0), Inches(5.5), Inches(4.3), Inches(0.7),
         [[("문턱이 낮아지면 검토가 늘고, 검토 하나하나가 매칭 풀이 됩니다", 12.5, REG, MUTED, 0)]], ls=1.4)
    footer(s, "10")

    # ===== S11 수익 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "발전 가능성 · 매출", "매출이 발생하는 경로")
    stages = [("1단계", "B2G 지자체 구독", "승계 모니터링 대시보드 · 3개 도 48개 시군 1차 시장", GREEN),
              ("2단계", "성사 성공보수 · 중개사 멤버십", "거래는 제휴 공인중개사, 팜바톤은 플랫폼 연계 대가", RGBColor(0x7A, 0x8A, 0x66)),
              ("상시", "정밀 리포트 · 제휴", "실사 연계 A·B등급 리포트(유료) · 정책자금·보험 제휴", SAGE)]
    sy = Inches(2.35)
    for i, (tag, t, d, col) in enumerate(stages):
        y = Emu(int(sy) + i * int(Inches(1.12)))
        text(s, LM, y, Inches(1.1), Inches(0.8), [[(tag, 13, XBOLD, col, 0)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.95), y, Inches(4.5), Inches(0.9),
             [[(t, 15, SEMI, INK, 0)], [(d, 11.5, REG, MUTED, 0)]], ls=1.35, sa=2)
        hline(s, LM, Emu(int(y) + int(Inches(0.98))), Inches(5.6), HAIR, 0.75)
    text(s, LM, Inches(5.85), Inches(5.6), Inches(0.4),
         [[("개인(농장주·청년농)은 무료", 13, SEMI, GREEN_DK, 0), (" — 매칭 풀이 플랫폼의 생명", 12, REG, MUTED, 0)]])
    pw = Inches(6.35); ph = Emu(int(pw * 9 / 16))
    s.shapes.add_picture(str(assets / "slide11@2x.png"), Inches(6.6), Inches(2.3), pw, ph)
    text(s, Inches(6.6), Emu(int(Inches(2.3)) + int(ph) + int(Inches(0.06))), pw, Inches(0.3),
         [[("지자체 구독 대시보드 — 콘셉트 목업(예시 데이터)", 10, REG, FAINT, 0)]], align=PP_ALIGN.CENTER)
    hline(s, LM, Inches(6.4), CW, HAIR, 1.0)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.62), CW, Inches(0.5),
         [[("3개 도 과수원 실거래 5,122건 · 중앙값 5,958만 원", 15, XBOLD, INK, -0.3),
           ("   저희 DB에서 실측한 숫자입니다", 12.5, REG, MUTED, 0)]])
    footer(s, "11")

    # ===== S12 확장 로드맵 (신규 2026-07-20 — 시간 미검증, 단축 시 1순위 컷) =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "발전 가능성 · 확장", "지금 구조 그대로, 세 방향으로 넓어집니다")
    expansions = [("작목", "품종 확장", "수령계수·소득계수 등 기준 데이터만 추가 — 산식·코드는 그대로", GREEN),
                  ("지역", "전국 확장", "팜맵·공시지가·KAMIS는 전국 단위 데이터 — 대상 지역만 넓히면 전국", RGBColor(0x7A, 0x8A, 0x66)),
                  ("정주여건", "매칭 요인 확장", "교육·의료·문화체육·교통 공공데이터를 매칭 점수에 결합 (5~6종 추가)", SAGE)]
    ey = Inches(2.35)
    for i, (tag, t, d, col) in enumerate(expansions):
        y = Emu(int(ey) + i * int(Inches(1.15)))
        text(s, LM, y, Inches(1.6), Inches(0.8), [[(tag, 13, XBOLD, col, 0)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(2.2), y, Inches(9.3), Inches(0.9),
             [[(t, 16, SEMI, INK, 0)], [(d, 12.5, REG, MUTED, 0)]], ls=1.35, sa=2)
        hline(s, LM, Emu(int(y) + int(Inches(1.0))), CW, HAIR, 0.75)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.15), CW, Inches(0.5),
         [[("전부 지금의 산식·매칭 엔진 구조를 바꾸지 않고, 데이터만 추가하는 확장입니다", 15, XBOLD, GREEN_DK, -0.3)]])
    footer(s, "12")

    # ===== S13 팀 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "팀", "개발 + 디자인 + 농업 도메인")
    team = [("김정훈", "대표 · 컴퓨터공학", ["기획·풀스택 개발 총괄", "공공데이터 적재 · 가치평가·매칭 엔진 · 배포"]),
            ("윤채원", "컴퓨터공학", ["UI/UX 디자인 · 사용성 검토", "모바일 퍼스트 화면 설계"]),
            ("유수민", "농업 계열", ["농업 도메인 자문", "승계 현장 이해 · 작목·시설 검증"])]
    x0, colw = LM, Inches(3.9)
    for i, (name, role, desc) in enumerate(team):
        x = Emu(int(x0) + i * int(colw))
        if i > 0:
            vline(s, Emu(int(x) - Inches(0.25)), Inches(2.8), Inches(2.3), HAIR, 1.0)
        text(s, x, Inches(2.75), Inches(3.5), Inches(0.7), [[(name, 30, BLACK, FOREST, -0.5)]])
        text(s, x, Inches(3.65), Inches(3.5), Inches(0.35), [[(role, 13, SEMI, GREEN_DK, 0.3)]])
        text(s, x, Inches(4.2), Inches(3.5), Inches(1.0),
             [[(d, 12.5, REG, MUTED, 0)] for d in desc], ls=1.4, sa=3)
    footer(s, "13")

    # ===== S14 클로징 =====
    s = new(); bg(s, FOREST)
    text(s, LM, Inches(1.9), Inches(11.5), Inches(0.5), [[("승계 검토의 진입장벽을", 26, SEMI, PALE, -0.3)]])
    text(s, Emu(int(LM) - Inches(0.02)), Inches(2.6), Inches(12), Inches(1.1),
         [[("수일 · 수십만 원", 40, XBOLD, WHITE, -1.0), ("   →   ", 30, MED, RGBColor(0x5A, 0x6E, 0x60), 0),
           ("무료 · 수십 초", 46, BLACK, SAGE, -1.5)]])
    text(s, Emu(int(LM) + Inches(0.02)), Inches(4.05), Inches(11.5), Inches(0.5),
         [[("사라지던 농장을 데이터로 다음 세대에 잇겠습니다", 18, REG, PALE, 0)]])
    hline(s, Emu(int(LM) + Inches(0.02)), Inches(5.35), Inches(11.6), HAIR_DK, 1.0)
    text(s, LM, Inches(5.6), Inches(8), Inches(0.5),
         [[("farmbaton.vercel.app", 18, SEMI, WHITE, 0), ("   지금 접속해 보세요", 13, REG, RGBColor(0x8F, 0xA5, 0x95), 0)]])
    logo_mark(s, LM, Inches(6.87), on_dark=True)
    text(s, Emu(int(LM) + int(MARK_SIZE) + Inches(0.1)), Inches(6.9), Inches(9), Inches(0.3),
         [[("감사합니다 · 김정훈 · 윤채원 · 유수민", 12, MED, RGBColor(0x6E, 0x78, 0x6F), 0.5)]])
    # QR
    block(s, Inches(11.0), Inches(5.5), Inches(1.5), Inches(1.5), WHITE, radius=0.06)
    s.shapes.add_picture(str(qr), Inches(11.14), Inches(5.64), Inches(1.22), Inches(1.22))

    prs.save(out)
    print(f"saved: {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    assets_dir = Path(sys.argv[1] if len(sys.argv) > 1 else ".")
    out_path = Path(__file__).resolve().parent.parent / "docs" / "발표_슬라이드.pptx"
    build(assets_dir, out_path)
