# -*- coding: utf-8 -*-
"""팜바톤 발표 템플릿(.pptx) 생성 — 복제해서 내용만 바꿔 쓰는 레이아웃 패턴 10종.

근거: docs/발표_디자인시스템.md (토큰·좌표·컴포넌트 스펙의 단일 진실 원천)
헬퍼·색상은 scripts/build_presentation.py 를 그대로 재사용한다.

사용법: python scripts/build_template.py
출력: docs/발표_템플릿.pptx
"""
import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

_HERE = Path(__file__).resolve().parent
_spec = importlib.util.spec_from_file_location("bp", str(_HERE / "build_presentation.py"))
bp = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(bp)

# 토큰 별칭 (디자인 시스템 문서의 이름 ↔ 구현 상수)
PAPER, FOREST, GREEN, GREEN_DK = bp.PAPER, bp.FOREST, bp.GREEN, bp.GREEN_DK
SAGE, SAGE_DK, INK, MUTED, FAINT = bp.SAGE, bp.SAGE_DK, bp.INK, bp.MUTED, bp.FAINT
HAIR, HAIR_DK, WHITE, PALE, AMBER, TINT = bp.HAIR, bp.HAIR_DK, bp.WHITE, bp.PALE, bp.AMBER, bp.TINT
BLACK, XBOLD, SEMI, MED, REG = bp.BLACK, bp.XBOLD, bp.SEMI, bp.MED, bp.REG
SW, SH, LM, CW = bp.SW, bp.SH, bp.LM, bp.CW
text, hline, vline, block = bp.text, bp.hline, bp.vline, bp.block
bg, kicker_head, footer, logo_mark = bp.bg, bp.kicker_head, bp.footer, bp.logo_mark


def build(out: Path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    def new():
        return prs.slides.add_slide(blank)

    def image_placeholder(s, x, y, w, h, label="이미지 · 목업 영역"):
        sp = block(s, x, y, w, h, WHITE, radius=0.04)
        sp.line.color.rgb = HAIR
        sp.line.width = bp.Pt(1.0)
        text(s, x, y, w, h, [[(label, 13, MED, FAINT, 0.5)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ===== T0 사용 안내 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "팜바톤 발표 템플릿", "복제해서 내용만 바꿔 쓰세요")
    guide = [
        ("T1  다크 타이틀", "발표 시작 — 다크 밴드 + 브랜드 카피"),
        ("T2  기본 콘텐츠", "킥커·헤드라인 + 본문 + 하단 결론 한 줄"),
        ("T3  빅넘버 3열", "핵심 수치 — 위기 지표 열만 앰버"),
        ("T4  비교표", "헤어라인 표 — 팜바톤 열만 틴트 면"),
        ("T5  단계 리스트", "수익모델·로드맵형 행 리스트"),
        ("T6  좌우 분할", "좌 텍스트 스택 + 우 이미지·목업"),
        ("T7  3열 카드", "팀 소개 등 동급 3열"),
        ("T8  차트·다이어그램", "풀 영역 + 출처 캡션 자리"),
        ("T9  다크 클로징", "발표 끝 — 각인 카피 + URL·QR"),
    ]
    col_h = Inches(0.62)
    for i, (t, d) in enumerate(guide):
        col, row = divmod(i, 5)
        x = Emu(int(LM) + col * int(Inches(6.0)))
        y = Emu(int(Inches(2.3)) + row * int(col_h))
        text(s, x, y, Inches(2.3), Inches(0.5), [[(t, 14, SEMI, INK, 0)]])
        text(s, Emu(int(x) + int(Inches(2.35))), y, Inches(3.5), Inches(0.5),
             [[(d, 12, REG, MUTED, 0)]])
    hline(s, LM, Inches(5.75), CW, HAIR, 1.0)
    text(s, LM, Inches(5.95), CW, Inches(0.9),
         [[("규칙 요약 — 자세한 것은 docs/발표_디자인시스템.md", 12.5, SEMI, GREEN_DK, 0)],
          [("헤드라인·숫자는 Black/ExtraBold, 나머지는 SemiBold 이하(700 금지) · 초록 강조는 화면당 1회 · 위기 수치는 앰버 전담", 11.5, REG, MUTED, 0)],
          [("구획은 헤어라인, 그림자·그라데이션·이모지 금지 · 라운드 14px 초과 금지 · 문단 텍스트 금지(대사로)", 11.5, REG, MUTED, 0)]],
         ls=1.5, sa=4)
    footer(s, "T0")

    # ===== T1 다크 타이틀 =====
    s = new(); bg(s, FOREST)
    text(s, LM, Inches(0.72), Inches(11.6), Inches(0.3),
         [[("행사·대회명 (킥커 자리)", 11.5, MED, SAGE, 2.2)]])
    hline(s, Emu(int(LM) + Inches(0.02)), Inches(1.16), Inches(3.2), HAIR_DK, 1.0)
    text(s, Inches(0.8), Inches(2.3), Inches(11), Inches(2.0), [[("타이틀", 96, BLACK, WHITE, -2.0)]])
    text(s, Inches(0.9), Inches(4.12), Inches(11), Inches(0.5), [[("Sub Title", 22, MED, PALE, 4.0)]])
    text(s, Emu(int(LM) + Inches(0.02)), Inches(5.12), Inches(11.5), Inches(0.7),
         [[("한 줄 카피 — 강조어는 ", 27, XBOLD, WHITE, -0.5), ("세이지 색 런", 27, BLACK, SAGE, -0.5)]])
    hline(s, Emu(int(LM) + Inches(0.02)), Inches(6.5), Inches(11.6), HAIR_DK, 1.0)
    logo_mark(s, LM, Inches(6.69), on_dark=True)
    text(s, Emu(int(LM) + int(bp.MARK_SIZE) + Inches(0.1)), Inches(6.72), Inches(8), Inches(0.4),
         [[("url-or-key-fact", 15, SEMI, WHITE, 0), ("   보조 문구", 12.5, REG, RGBColor(0x8F, 0xA5, 0x95), 0)]])
    text(s, Inches(9.5), Inches(6.72), Inches(2.95), Inches(0.4),
         [[("발표자 이름들", 11.5, MED, RGBColor(0x8F, 0xA5, 0x95), 0.5)]], align=PP_ALIGN.RIGHT)

    # ===== T2 기본 콘텐츠 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "카테고리 · 킥커", "헤드라인 — 슬라이드당 메시지 하나")
    text(s, LM, Inches(2.5), Inches(11.0), Inches(3.0),
         [[("본문 영역", 15, SEMI, INK, 0)],
          [("문단을 쓰지 말고 구·라벨 단위로. 서술은 발표 대사에 둔다.", 12.5, REG, MUTED, 0)],
          [("행이 여러 개면 행 사이에 헤어라인.", 12.5, REG, MUTED, 0)]],
         ls=1.5, sa=6)
    hline(s, LM, Inches(6.12), CW, HAIR, 1.0)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.36), Inches(11.6), Inches(0.6),
         [[("하단 결론 한 문장 — ", 19, MED, INK, -0.3), ("강조 런은 초록", 19, XBOLD, GREEN, -0.3)]])
    footer(s, "T2")

    # ===== T3 빅넘버 3열 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "카테고리 · 킥커", "빅넘버 헤드라인")
    cols = [("00.0", "%", "지표 라벨", "부연 설명 한 줄", FOREST),
            ("00", "만 가구", "지표 라벨", "부연 설명 한 줄", FOREST),
            ("0.0", "%", "위기 지표 라벨", "위기·감소 지표만 앰버", AMBER)]
    x0, colw, gap, top = LM, Inches(3.85), Inches(0.15), Inches(2.7)
    for i, (big, unit, lbl, sub, col) in enumerate(cols):
        x = Emu(int(x0) + i * (int(colw) + int(gap)))
        if i > 0:
            vline(s, Emu(int(x) - int(gap) // 2), Inches(2.82), Inches(2.35), HAIR, 1.0)
        text(s, x, top, colw, Inches(1.3), [[(big, 74, BLACK, col, -2.0), (unit, 22, MED, MUTED, 0)]])
        hline(s, Emu(int(x) + Inches(0.05)), Inches(4.28), Inches(2.4), col, 2.5)
        text(s, x, Inches(4.45), colw, Inches(0.9),
             [[(lbl, 15, SEMI, INK, 0)], [(sub, 11.5, REG, MUTED, 0)]], ls=1.25, sa=2)
    hline(s, LM, Inches(6.12), CW, HAIR, 1.0)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.36), Inches(11.6), Inches(0.6),
         [[("하단 결론 한 문장 — ", 19, MED, INK, -0.3), ("강조 런", 19, XBOLD, GREEN, -0.3)]])
    text(s, LM, Inches(6.78), Inches(8), Inches(0.3), [[("출처: 기관 · 조사명(연도)", 10, REG, FAINT, 0.5)]])
    footer(s, "T3")

    # ===== T4 비교표 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "카테고리 · 킥커", "비교표 헤드라인")
    tx, c0, cw_, fbw = LM, Inches(2.0), Inches(2.55), Inches(2.35)
    xs = [tx, Emu(int(tx) + int(c0)), Emu(int(tx) + int(c0) + int(cw_)), Emu(int(tx) + int(c0) + 2 * int(cw_))]
    fb_x = Emu(int(tx) + int(CW) - int(fbw))
    top = Inches(2.15)
    block(s, Emu(int(fb_x) - Inches(0.22)), top, Emu(int(fbw) + Inches(0.22)), Inches(4.15), TINT, radius=0.04)
    hline(s, Emu(int(fb_x) - Inches(0.22)), top, Emu(int(fbw) + Inches(0.22)), INK, 2.0)
    heads = ["", "비교 대상 A", "비교 대상 B", "비교 대상 C"]
    for x, h in zip(xs, heads):
        text(s, x, Inches(2.35), Inches(2.4), Inches(0.4), [[(h, 13, MED, MUTED, 0)]])
    text(s, Emu(int(fb_x) - Inches(0.02)), Inches(2.35), fbw, Inches(0.4), [[("팜바톤", 13.5, XBOLD, GREEN_DK, 0)]])
    hline(s, tx, Inches(2.82), CW, INK, 1.25)
    rows = [("기준 1", "값", "값", "값", "차별점"),
            ("기준 2", "값", "값", "값", "차별점"),
            ("기준 3", "값", "값", "값", "차별점")]
    ry = Inches(3.15)
    for i, (crit, a, b_, c, fb) in enumerate(rows):
        y = Emu(int(ry) + i * int(Inches(1.0)))
        text(s, xs[0], y, Inches(1.9), Inches(0.4), [[(crit, 13.5, SEMI, INK, 0)]])
        for x, v in zip(xs[1:], (a, b_, c)):
            text(s, x, y, Inches(2.4), Inches(0.4), [[(v, 13, REG, MUTED, 0)]])
        text(s, Emu(int(fb_x) - Inches(0.02)), y, fbw, Inches(0.4), [[(fb, 13.5, SEMI, INK, 0)]])
        if i < 2:
            hline(s, tx, Emu(int(y) + int(Inches(0.72))), CW, HAIR, 0.75)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.7), Inches(11.7), Inches(0.5),
         [[("표 아래 결론 한 문장 ", 16, XBOLD, INK, -0.3), (" 보조 설명", 13, REG, MUTED, 0)]])
    footer(s, "T4")

    # ===== T5 단계 리스트 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "카테고리 · 킥커", "단계 리스트 헤드라인")
    stages = [("1단계", "단계 제목", "설명 한 줄 — caption 스타일", GREEN),
              ("2단계", "단계 제목", "설명 한 줄", SAGE_DK),
              ("상시", "단계 제목", "설명 한 줄", SAGE)]
    sy = Inches(2.35)
    for i, (tag, t, d, col) in enumerate(stages):
        y = Emu(int(sy) + i * int(Inches(1.15)))
        text(s, LM, y, Inches(1.6), Inches(0.8), [[(tag, 13, XBOLD, col, 0)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(2.2), y, Inches(9.3), Inches(0.9),
             [[(t, 16, SEMI, INK, 0)], [(d, 12.5, REG, MUTED, 0)]], ls=1.35, sa=2)
        hline(s, LM, Emu(int(y) + int(Inches(1.0))), CW, HAIR, 0.75)
    text(s, Emu(int(LM) - Inches(0.02)), Inches(6.15), CW, Inches(0.5),
         [[("리스트 아래 결론 한 문장 — 초록 강조", 15, XBOLD, GREEN_DK, -0.3)]])
    footer(s, "T5")

    # ===== T6 좌우 분할 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "카테고리 · 킥커", "좌우 분할 헤드라인")
    items = [("포인트 제목", "설명 한 줄"), ("포인트 제목", "설명 한 줄"),
             ("포인트 제목", "설명 한 줄"), ("포인트 제목", "설명 한 줄")]
    bx, bw_ = LM, Inches(4.4)
    for i, (t, d) in enumerate(items):
        y = Emu(int(Inches(2.35)) + i * int(Inches(1.02)))
        text(s, bx, y, bw_, Inches(0.9),
             [[(t, 15, SEMI, INK, 0)], [(d, 11.5, REG, MUTED, 0)]], ls=1.35, sa=3)
        if i < 3:
            hline(s, bx, Emu(int(y) + int(Inches(0.82))), bw_, HAIR, 0.75)
    vline(s, Inches(5.5), Inches(2.45), Inches(3.6), HAIR, 1.0)
    image_placeholder(s, Inches(5.95), Inches(2.35), Inches(6.5), Inches(3.66))
    text(s, Inches(5.95), Inches(6.08), Inches(6.5), Inches(0.3),
         [[("이미지 캡션 — 목업이면 '콘셉트 목업(예시 데이터)' 표기", 10, REG, FAINT, 0)]], align=PP_ALIGN.CENTER)
    text(s, LM, Inches(6.75), CW, Inches(0.4),
         [[("하단 보조 정보 한 줄 (fine · mute)", 11.5, REG, FAINT, 0)]])
    footer(s, "T6")

    # ===== T7 3열 카드 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "카테고리 · 킥커", "3열 카드 헤드라인")
    team = [("이름", "역할 · 전공", ["담당 영역 한 줄", "담당 영역 한 줄"]),
            ("이름", "역할 · 전공", ["담당 영역 한 줄", "담당 영역 한 줄"]),
            ("이름", "역할 · 전공", ["담당 영역 한 줄", "담당 영역 한 줄"])]
    x0, colw = LM, Inches(3.9)
    for i, (name, role, desc) in enumerate(team):
        x = Emu(int(x0) + i * int(colw))
        if i > 0:
            vline(s, Emu(int(x) - Inches(0.25)), Inches(2.8), Inches(2.3), HAIR, 1.0)
        text(s, x, Inches(2.75), Inches(3.5), Inches(0.7), [[(name, 30, BLACK, FOREST, -0.5)]])
        text(s, x, Inches(3.65), Inches(3.5), Inches(0.35), [[(role, 13, SEMI, GREEN_DK, 0.3)]])
        text(s, x, Inches(4.2), Inches(3.5), Inches(1.0),
             [[(d, 12.5, REG, MUTED, 0)] for d in desc], ls=1.4, sa=3)
    footer(s, "T7")

    # ===== T8 차트·다이어그램 =====
    s = new(); bg(s, PAPER)
    kicker_head(s, "카테고리 · 킥커", "차트·다이어그램 헤드라인")
    image_placeholder(s, LM, Inches(2.2), CW, Inches(4.2),
                      "차트 · 다이어그램 영역  (선 2.5px · 직접 라벨은 처음·끝만 · 위기 데이터는 앰버)")
    text(s, LM, Inches(6.6), Inches(8), Inches(0.3),
         [[("출처: 기관 · 조사명(연도)", 10, REG, FAINT, 0.5)]])
    footer(s, "T8")

    # ===== T9 다크 클로징 =====
    s = new(); bg(s, FOREST)
    text(s, LM, Inches(1.9), Inches(11.5), Inches(0.5), [[("리드 문장 (PALE 색)", 26, SEMI, PALE, -0.3)]])
    text(s, Emu(int(LM) - Inches(0.02)), Inches(2.6), Inches(12), Inches(1.1),
         [[("비포 카피", 40, XBOLD, WHITE, -1.0), ("   →   ", 30, MED, RGBColor(0x5A, 0x6E, 0x60), 0),
           ("애프터 카피", 46, BLACK, SAGE, -1.5)]])
    text(s, Emu(int(LM) + Inches(0.02)), Inches(4.05), Inches(11.5), Inches(0.5),
         [[("마무리 문장 한 줄", 18, REG, PALE, 0)]])
    hline(s, Emu(int(LM) + Inches(0.02)), Inches(5.35), Inches(11.6), HAIR_DK, 1.0)
    text(s, LM, Inches(5.6), Inches(8), Inches(0.5),
         [[("url-or-cta", 18, SEMI, WHITE, 0), ("   보조 문구", 13, REG, RGBColor(0x8F, 0xA5, 0x95), 0)]])
    qr = block(s, Inches(11.0), Inches(5.5), Inches(1.5), Inches(1.5), WHITE, radius=0.06)
    text(s, Inches(11.0), Inches(5.5), Inches(1.5), Inches(1.5), [[("QR", 16, SEMI, FAINT, 1.0)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    logo_mark(s, LM, Inches(6.87), on_dark=True)
    text(s, Emu(int(LM) + int(bp.MARK_SIZE) + Inches(0.1)), Inches(6.9), Inches(9), Inches(0.3),
         [[("감사 인사 · 발표자 이름들", 12, MED, RGBColor(0x6E, 0x78, 0x6F), 0.5)]])

    prs.save(str(out))
    print(f"saved: {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build(Path(__file__).resolve().parent.parent / "docs" / "발표_템플릿.pptx")
